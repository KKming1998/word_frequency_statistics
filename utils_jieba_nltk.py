import os
import re
import jieba
from collections import Counter
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd
import PyPDF2
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
import argparse
import pandas as pd
from datetime import datetime
import logging
import comtypes.client



# 配置日志
logging.basicConfig(filename='word_count.log', level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')

# 确保下载 punkt 和 stopwords 资源
#nltk.download('punkt')
#nltk.download('stopwords')

nltk_data_path = os.path.join(os.getcwd(), 'nltk_data')
nltk.data.path.append(nltk_data_path)

# 定义常见的中文语气词列表
chinese_stopwords = set([
    "啊", "吧", "吗", "呢", "哦", "嗯", "呀", "了", "的", "地", "得", "着", "在", "上", "下", "里", "中", "间",
    "前", "后", "内", "外", "下", "到", "对", "和", "跟", "同", "如", "与", "及", "而", "或", "即", "也", "就",
    "是", "被", "把", "被", "让", "给", "为", "以", "因", "由", "从", "自", "向", "往", "当", "若", "如", "若",
    "怎", "何", "谁", "哪", "几", "多", "有", "无", "不", "没", "别", "再", "可", "要", "会", "能", "可", "该",
    "这", "那", "些", "些", "此", "其", "彼", "另", "各", "每", "两", "三", "四", "五", "六", "七", "八", "九",
    "十", "百", "千", "万", "亿", "一", "二", "两", "三", "四", "五", "六", "七", "八", "九", "十"
])

# 加载英语停用词列表
english_stopwords = set(stopwords.words('english'))

def read_blacklist(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return set(word.strip() for word in file.readlines())

def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


def read_doc(file_path):

    # 获取文件的绝对路径
    absolute_path = os.path.abspath(file_path)

    # 初始化 Word 应用程序对象
    word = comtypes.client.CreateObject('Word.Application')
    word.Visible = False  # 不显示 Word 程序窗口

    try:
        # 打开指定路径的 .doc 文件
        doc = word.Documents.Open(absolute_path)

        # 读取文档内容
        text = doc.Content.Text

        # 关闭文档
        doc.Close(False)  # 不保存更改
    finally:
        # 退出 Word 应用程序
        word.Quit()

    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text



def read_ppt(ppt_file):
    """
    读取.ppt文件的内容，并将其合并为一个字符串。

    :param ppt_file: .ppt文件的路径
    :return: 包含所有幻灯片内容的字符串
    """
    # 获取文件的绝对路径
    absolute_path = os.path.abspath(ppt_file)

    # 创建PowerPoint应用程序对象
    powerpoint = comtypes.client.CreateObject('PowerPoint.Application')
    powerpoint.Visible = True  # 不显示PowerPoint应用程序窗口

    try:
        # 打开.ppt文件
        presentation = powerpoint.Presentations.Open(absolute_path)

        # 初始化一个空字符串来存储所有幻灯片的内容
        all_slides_content = ""

        # 读取每个幻灯片的内容
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text_frame = shape.TextFrame
                    if text_frame.HasText:
                        text_range = text_frame.TextRange
                        all_slides_content += text_range.Text + "\n"  # 将文本添加到字符串末尾，并换行

        # 关闭演示文稿
        presentation.Close()
    except Exception as e:
        print(f"读取文件时发生错误: {e}")
        all_slides_content = ""
    finally:
        # 退出PowerPoint应用程序
        powerpoint.Quit()

    return all_slides_content

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def read_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    text = []
    for sheet in wb.sheetnames:
        for row in wb[sheet].iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text.append(str(cell))
    return ' '.join(text)

def read_xls(file_path):
    wb = xlrd.open_workbook(file_path)
    text = []
    for sheet in wb.sheets():
        for row_idx in range(sheet.nrows):
            for col_idx in range(sheet.ncols):
                value = sheet.cell(row_idx, col_idx).value
                if value:
                    text.append(str(value))
    return ' '.join(text)

def read_xlsm(file_path):
    # 使用 openpyxl 读取 .xlsm 文件
    return read_xlsx(file_path)

def read_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    return text

def clean_text(text):
    # 去除标点符号
    text = re.sub(r'[^\w\s]', '', text)
    # 去除数字
    text = re.sub(r'\d+', '', text)
    return text

def word_count(text, blacklist):
    # 分离中文和英文
    chinese_text = re.findall(r'[\u4e00-\u9fff]+', text)
    english_text = re.findall(r'[a-zA-Z]+', text)

    # 中文分词
    chinese_words = []
    for segment in chinese_text:
        chinese_words.extend(jieba.lcut(segment))

    # 英文分词
    english_words = word_tokenize(' '.join(english_text))

    # 合并中英文词汇
    all_words = chinese_words + english_words

    # 过滤掉停用词和黑名单词
    filtered_words = [word for word in all_words if
                      word not in chinese_stopwords and word not in english_stopwords and word not in blacklist and len(word) > 1]

    return dict(Counter(filtered_words))

def get_top_n_words(word_counts, n=10):
    # 对词频字典按值降序排序
    sorted_word_counts = sorted(word_counts.items(), key=lambda x: x[1], reverse=True)
    # 取前n个
    return dict(sorted_word_counts[:n])

def process_file(file_path, blacklist, top_n):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.txt':
        text = read_txt(file_path)
    elif ext == '.doc':
        text = read_doc(file_path)
    elif ext == '.docx':
        text = read_docx(file_path)
    elif ext == '.ppt':
        text = read_ppt(file_path)
    elif ext in ['.pptx', '.pptm']:
        text = read_pptx(file_path)
    elif ext in ['.xlsx', '.xlsm']:
        text = read_xlsx(file_path)
    elif ext == '.xls':
        text = read_xls(file_path)
    elif ext == '.pdf':
        text = read_pdf(file_path)
    else:
        logging.warning(f"Unsupported file type: {file_path}")
        return None

    cleaned_text = clean_text(text)
    word_counts = word_count(cleaned_text, blacklist)
    top_words = get_top_n_words(word_counts, n=top_n)
    logging.info(f"Top {top_n} words in {file_path}: {top_words}")
    return top_words

def merge_word_counts(word_counts_list):
    merged_counts = Counter()
    for word_counts in word_counts_list:
        if word_counts:
            merged_counts.update(word_counts)
    return dict(merged_counts)

def process_directory(directory, blacklist, top_n):
    results = []
    unsupported_files = []
    processed_files = 0

    for root, dirs, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            try:
                word_counts = process_file(file_path, blacklist, top_n)
                if word_counts:
                    results.append(word_counts)
                    processed_files += 1
                else:
                    unsupported_files.append(file_path)
            except Exception as e:
                logging.error(f"Error processing {file_path}: {e}")
                unsupported_files.append(file_path)

    logging.info(f"Total processed files: {processed_files}")
    logging.info(f"Unsupported files: {unsupported_files}")
    return merge_word_counts(results), processed_files, unsupported_files

def sort_word_counts(word_counts):
    return dict(sorted(word_counts.items(), key=lambda x: x[1], reverse=True))

def scale_and_cut_word_counts(word_counts, weight_coefficient, cut_n):
    scaled_word_counts = {word: round(count * weight_coefficient) for word, count in word_counts.items()}
    sorted_scaled_word_counts = sorted(scaled_word_counts.items(), key=lambda x: x[1], reverse=True)
    return dict(sorted_scaled_word_counts[:cut_n])

def write_to_csv(word_counts, directory, current_time):
    df = pd.DataFrame(word_counts.items(), columns=['Word', 'Count'])
    output_file = f"{os.path.basename(directory)}_{current_time}.csv"
    df.to_csv(output_file, index=False)
    logging.info(f"Results written to {output_file}")

def main():
    parser = argparse.ArgumentParser(description="Process files and output top N most frequent words.")
    parser.add_argument('directory', help='Directory to process')
    parser.add_argument('--blacklist', help='Path to a file containing blacklist words, one per line')
    parser.add_argument('--top-n', type=int, default=10, help='Number of top words to keep (default: 10)')
    parser.add_argument('--weight-coefficient', type=float, default=1.0, help='Weight coefficient for scaling (default: 1.0)')
    parser.add_argument('--cut-n', type=int, default=10, help='Number of words to keep after scaling and cutting (default: 10)')
    args = parser.parse_args()

    start_time = datetime.now()
    logging.info(f"Program started at {start_time}")

    blacklist = set()
    if args.blacklist:
        blacklist = read_blacklist(args.blacklist)

    word_counts, processed_files, unsupported_files = process_directory(args.directory, blacklist, args.top_n)
    sorted_word_counts = sort_word_counts(word_counts)
    scaled_cut_word_counts = scale_and_cut_word_counts(sorted_word_counts, args.weight_coefficient, args.cut_n)

    current_time = datetime.now().strftime("%Y%m%d_%H%M%S")
    write_to_csv(scaled_cut_word_counts, args.directory, current_time)

    end_time = datetime.now()
    logging.info(f"Program ended at {end_time}")
    logging.info(f"Total time taken: {end_time - start_time}")

if __name__ == "__main__":
    main()